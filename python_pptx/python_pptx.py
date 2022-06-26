#!/usr/bin/env python

import os
import numpy as np
import cv2
from pptx import Presentation

class MyPresentation():
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)

        self.prs = Presentation()

        #Default slide size is
        # 9144000 x 6858000 (4:3)
        # 12193200 x 6858000(16:9)
        self.prs.slide_width = 12193200
        self.prs.slide_height = 6858000

    def add_blank_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        return slide

    def add_images_to_region(self, slide_idx, t, b, l, r, fnames_img, titles=None,
                             align="vertical"):
        """
        slide_idx:  index of slide in this presentation
                      (Don't forget to add one before calling this function)
        t, b, l, r: top, bottom, left, right boundary of the region (in % of slide size)
        fnames_img: list of image files
        titles:     image titles
        align:      "vertical" or "horizontal"
        """
        assert slide_idx < len(self.prs.slides), \
            f"Index {slide_idx}/{len(self.prs.slides)} out of range"

        img_shapes = np.array([(cv2.imread(f)).shape[:2] for f in fnames_img],
                              dtype=np.float32)

        if align == "vertical":
            # Assuming height = 1, then
            resized_width_ratio = img_shapes[:, 1] / img_shapes[:, 0] # image aspect ratio
            combined_width_ratio = np.sum(resized_width_ratio) # combined aspect ratio

            region_aspect = (r - l)*self.prs.slide_width / ((b - t)*self.prs.slide_height)
            if combined_width_ratio > region_aspect:
                region_width = self.prs.slide_width * (r - l)
                region_height = region_width / combined_width_ratio
            else:
                region_height = self.prs.slide_height * (b - t)
                region_width = region_height * combined_width_ratio

            region_left = self.prs.slide_width * (l + r)/2 - region_width/2
            region_top = self.prs.slide_height * (t + b)/2 - region_height/2

            img_height = region_height
            left = region_left
            for i, fname in enumerate(fnames_img):
                img_width = resized_width_ratio[i] * img_height
                self.prs.slides[slide_idx].shapes.add_picture(fname, left, region_top,
                                                              img_width, img_height)
                if titles is not None:
                    text_height = self.prs.slide_height*0.05
                    tbox = self.prs.slides[0].shapes.add_textbox(left,
                                                                 region_top - text_height,
                                                                 img_width, text_height)
                    tbox.text_frame.text = titles[i]
                    
                left += img_width
            
        elif align == "horizontal":
            assert False, "Not implemented yet!"
        else:
            assert False, "Invalid align"

    def save(self, fname):
        self.prs.save(fname)

def test():
    prs = MyPresentation()
    prs.add_blank_slide()
    prs.add_images_to_region(0, 0.5, 1.0, 0.5, 1.0, ["im_2.bmp", "im_3.bmp"],
                             ["Airplane", "Dog"], "vertical")
    prs.save("test.pptx")

if __name__ == "__main__":
    test()


